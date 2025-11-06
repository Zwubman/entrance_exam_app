import io
import os
import uuid
import textwrap
from PIL import Image
from pptx import Presentation
from app.config.setting import settings

def extract_ppt_data(file_bytes: bytes):
    file_stream = io.BytesIO(file_bytes)
    extracted_payloads = []

    try:
        prs = Presentation(file_stream)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            page_texts, page_images = [], []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    sentences = textwrap.wrap(shape.text.strip(), width=settings.CHUNK_LENGTH)
                    page_texts.extend(sentences)

            # Extract images
            for shape in slide.shapes:
                if hasattr(shape, 'image') and shape.shape_type == 13:  # 13 is picture shape
                    try:
                        image_bytes = shape.image.blob
                        img = Image.open(io.BytesIO(image_bytes))
                        image_name = f"slide_{slide_num}_{uuid.uuid4().hex}.png"
                        image_path = os.path.join(settings.UPLOADS_DIR, image_name)
                        img.save(image_path, "PNG")
                        page_images.append(image_path)
                    except Exception as e:
                        print(f"Failed to extract image from slide {slide_num}: {e}")
                        continue

            extracted_payloads.append({
                "texts": page_texts,
                "images": page_images
            })

    except Exception:
        # Fallback for old .ppt format
        try:
            import textract
            file_stream.seek(0)
            text = textract.process(file_stream, extension='ppt').decode('utf-8')
            text_segments = []
            
            if text and text.strip():
                sentences = textwrap.wrap(text.strip(), width=settings.CHUNK_LENGTH)
                text_segments.extend(sentences)

            extracted_payloads.append({
                "texts": text_segments,
                "images": []
            })
        except ImportError:
            print("textract not installed for .ppt file support")
            extracted_payloads.append({"texts": [], "images": []})
        except Exception as e:
            print(f"Error processing with textract: {e}")
            extracted_payloads.append({"texts": [], "images": []})

    return extracted_payloads